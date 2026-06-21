# -*- coding: utf-8 -*-
"""SauceDemo case-study deck built from the real team artifacts:
   SauceDemo (Test Cases).xlsx  -> 99 cases / 82 PASS / 17 FAIL across 5 modules
   SauceDemo (Bug Report).xlsx  -> 17 open bugs (13 low / 4 medium severity)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x13, 0x2B, 0x4D)
GREEN  = RGBColor(0x3D, 0xDC, 0x84)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x5A, 0x63, 0x72)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF8)
RED    = RGBColor(0xD9, 0x3A, 0x3A)
AMBER  = RGBColor(0xE0, 0x8A, 0x1E)
PASSG  = RGBColor(0x1E, 0x8E, 0x3E)
ROWALT = RGBColor(0xE9, 0xED, 0xF4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(1.05), NAVY)
    rect(s, 0, Inches(1.05), SW, Pt(4), GREEN)
    text(s, Inches(0.55), Inches(0.1), Inches(12), Inches(0.5),
         [[(title, 24, True, WHITE)]], anchor=MSO_ANCHOR.TOP)
    if kicker:
        text(s, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.35),
             [[(kicker, 11.5, False, GREEN)]])


def footer(s, n):
    text(s, Inches(0.5), Inches(7.05), Inches(10), Inches(0.35),
         [[("Pengujian Perangkat Lunak  •  Studi Kasus: SauceDemo  •  Sumber: Test Cases & Bug Report (.xlsx)", 9, False, GREY)]])
    text(s, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
         [[(str(n), 9, True, GREY)]], align=PP_ALIGN.RIGHT)


def table(s, x, y, w, col_w, rows, header_fill=NAVY, fontsize=9.5,
          row_h=Inches(0.3), align_left_long=True):
    widths = [int(w * f) for f in col_w]
    cury = y
    for ri, row in enumerate(rows):
        curx = x; is_head = ri == 0
        rh = Inches(0.34) if is_head else row_h
        for ci, cell in enumerate(row):
            cw = widths[ci]
            fill = header_fill if is_head else (WHITE if (ri % 2 == 1) else ROWALT)
            cellsp = rect(s, curx, cury, cw, rh, fill)
            color, bold = (WHITE if is_head else DARK), is_head
            up = str(cell).strip().upper()
            if not is_head and up in ("PASS", "FAIL", "BLOCKED", "OPENED", "HIGH", "MEDIUM",
                                      "MEDUM", "LOW", "CRITICAL"):
                if up == "PASS": color, bold = PASSG, True
                elif up in ("FAIL", "CRITICAL", "HIGH"): color, bold = RED, True
                elif up in ("BLOCKED", "OPENED", "MEDIUM", "MEDUM"): color, bold = AMBER, True
                else: color, bold = GREY, True
            tf = cellsp.text_frame; tf.word_wrap = True
            tf.margin_left = Pt(5); tf.margin_right = Pt(4)
            tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (is_head or len(str(cell)) <= 10) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(cell)
            r.font.size = Pt(fontsize); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "Calibri"
            curx += cw
        cury += rh
    return cury


# ---- data from the workbooks ----
MODULES = [  # name, total, pass, fail, fail TC ids
    ("Login",     44, 38, 6, "TC-LP-026 … 031 (account lockout)"),
    ("Home Page", 22, 19, 3, "TC-MP-001, 003, 013"),
    ("Cart",      12, 10, 2, "TC-CP-006, 008"),
    ("Checkout",  15, 10, 5, "TC-CKO-004, 007, 008, 009, 015"),
    ("Footer",     6,  5, 1, "TC-F-007"),
]
TOTAL, PASS, FAIL = 99, 82, 17

# ============================================================ 1 TITLE
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Pt(5), GREEN)
text(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.5),
     [[("STUDI KASUS — PENGUJIAN PERANGKAT LUNAK", 16, True, GREEN)]])
text(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(2.0),
     [[("Test Execution & Test Reporting", 42, True, WHITE)],
      [("Pengujian Manual Aplikasi E-Commerce ", 28, True, WHITE), ("SauceDemo", 28, True, GREEN)]],
     line_spacing=1.05)
text(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(1.6),
     [[("Target Uji:  https://www.saucedemo.com/", 15, False, WHITE)],
      [("99 Test Case  •  82 PASS / 17 FAIL  •  17 Bug Dilaporkan", 14, True, GREEN)],
      [("Sumber data: SauceDemo (Test Cases).xlsx  &  SauceDemo (Bug Report).xlsx", 12, False, RGBColor(0xB8,0xC4,0xD6))],
      [("Tim: Zaki & Galih", 12, False, RGBColor(0xB8,0xC4,0xD6))]],
     line_spacing=1.25, space_after=3)

# ============================================================ 2 AGENDA
s = slide(); header(s, "Agenda", "Alur dokumen pengujian")
items = [
    ("01", "Tentang Aplikasi", "Eksplorasi SauceDemo & fitur"),
    ("02", "Scope & Modul Uji", "5 modul: Login, Home, Cart, Checkout, Footer"),
    ("03", "Test Case Summary", "99 case — rekap per modul"),
    ("04", "Test Execution Result", "Pass rate & status per modul"),
    ("05", "Bug Report", "17 bug — severity & priority"),
    ("06", "Bug Kritis", "3 bug priority High"),
    ("07", "Analisis Kualitas", "Dampak, risiko, retesting"),
    ("08", "Release Recommendation", "GO / Conditional / NO-GO"),
]
x0, y0 = Inches(0.6), Inches(1.45)
cw, chh, gap = Inches(6.0), Inches(1.25), Inches(0.18)
for i, (num, t, d) in enumerate(items):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + Inches(0.25)); y = y0 + row * (chh + gap)
    rect(s, x, y, cw, chh, LIGHT); rect(s, x, y, Inches(0.12), chh, GREEN)
    text(s, x + Inches(0.3), y, Inches(1.1), chh, [[(num, 30, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(1.4), y, cw - Inches(1.5), chh,
         [[(t, 15, True, DARK)], [(d, 10.5, False, GREY)]], anchor=MSO_ANCHOR.MIDDLE, space_after=2)
footer(s, 2)

# ============================================================ 3 ABOUT
s = slide(); header(s, "Tentang Aplikasi: SauceDemo", "Eksplorasi & identifikasi fitur")
text(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.7),
     [[("SauceDemo (Swag Labs) adalah aplikasi web simulasi ", 13, False, DARK),
       ("e-commerce", 13, True, NAVY),
       (" dari Sauce Labs untuk latihan pengujian. Menyediakan beberapa akun yang memunculkan perilaku berbeda.", 13, False, DARK)]],
     line_spacing=1.15)
text(s, Inches(0.55), Inches(2.2), Inches(6), Inches(0.4), [[("Fitur Utama", 15, True, NAVY)]])
feats = ["Autentikasi (login / logout)", "Katalog produk + PDP",
         "Sorting (nama / harga)", "Keranjang & badge count",
         "Checkout 3 langkah", "Footer & social links"]
for i, f in enumerate(feats):
    y = Inches(2.65) + i * Inches(0.42)
    rect(s, Inches(0.6), y + Inches(0.06), Inches(0.12), Inches(0.12), GREEN)
    text(s, Inches(0.85), y, Inches(5.3), Inches(0.4), [[(f, 12, False, DARK)]])
text(s, Inches(6.9), Inches(2.2), Inches(6), Inches(0.4), [[("Akun Uji (password: secret_sauce)", 15, True, NAVY)]])
acc = [["Username", "Perilaku"],
       ["standard_user", "Normal / baseline"],
       ["locked_out_user", "Diblokir saat login"],
       ["problem_user", "Bug UI & fungsi"],
       ["performance_glitch_user", "Loading lambat"],
       ["error_user", "Error saat checkout"],
       ["visual_user", "Anomali visual"]]
table(s, Inches(6.9), Inches(2.65), Inches(6.0), [0.52, 0.48], acc, fontsize=10.5, row_h=Inches(0.33))
footer(s, 3)

# ============================================================ 4 SCOPE
s = slide(); header(s, "Scope & Strategi Pengujian", "5 modul diuji • pendekatan black-box")
def panel(px, title, color, lst):
    rect(s, px, Inches(1.45), Inches(6.0), Inches(0.5), color)
    text(s, px, Inches(1.5), Inches(6), Inches(0.4), [[(title, 15, True, WHITE)]], align=PP_ALIGN.CENTER)
    for i, it in enumerate(lst):
        y = Inches(2.15) + i * Inches(0.5)
        rect(s, px, y, Inches(6.0), Inches(0.42), LIGHT)
        text(s, px + Inches(0.2), y, Inches(5.7), Inches(0.42), [[("•  " + it, 12, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
panel(Inches(0.55), "IN-SCOPE (5 MODUL)", NAVY,
      ["Login & autentikasi (semua role)", "Home page: katalog, sorting, menu",
       "Cart: add/remove, badge, checkout", "Checkout: info, overview, complete", "Footer: social links, ToS/Privacy"])
panel(Inches(6.85), "OUT-OF-SCOPE", GREY,
      ["Load / performance testing", "Security / penetration testing",
       "Pengujian API / backend", "Otomasi end-to-end penuh"])
text(s, Inches(0.55), Inches(5.0), Inches(12.2), Inches(1.6),
     [[("Pendekatan: ", 13, True, NAVY),
       ("Black-box manual testing oleh tim (Zaki & Galih) lintas role & device, "
        "termasuk skenario positif, negatif, edge case, dan responsivitas.", 13, False, DARK)],
      [("Status dipakai: ", 13, True, NAVY), ("PASS  •  FAIL  •  (Bug: Opened)", 13, False, DARK)]],
     line_spacing=1.2, space_after=8)
footer(s, 4)

# ============================================================ 5 TEST CASE SUMMARY
s = slide(); header(s, "Test Case Summary per Modul", "99 test case • SauceDemo (Test Cases).xlsx")
rows = [["Modul", "Total", "PASS", "FAIL", "Test Case yang FAIL"]]
for nm, tot, p, f, ids in MODULES:
    rows.append([nm, str(tot), str(p), str(f), ids])
rows.append(["TOTAL", str(TOTAL), str(PASS), str(FAIL), "17 bug dilaporkan"])
table(s, Inches(0.5), Inches(1.45), Inches(12.4), [0.14, 0.09, 0.09, 0.09, 0.59], rows,
      fontsize=11.5, row_h=Inches(0.55))
text(s, Inches(0.5), Inches(5.6), Inches(12), Inches(0.5),
     [[("Pass Rate keseluruhan: ", 15, True, NAVY), ("82.8%  (82 / 99)", 15, True, PASSG)]])
text(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.7),
     [[("Catatan: ", 12, True, NAVY),
       ("Modul Login mendominasi jumlah case (44) karena diuji untuk 6 role. Kegagalan terbesar di "
        "Checkout (5) dan Login (6 — semua terkait fitur account-lockout yang tidak ada.)", 12, False, DARK)]],
     line_spacing=1.15)
footer(s, 5)

# ============================================================ 6 KEY TEST CASES
s = slide(); header(s, "Test Case Terpenting (PASS & FAIL)", "13 case prioritas — alur inti + bug utama")
KEY = [
    # id, modul, skenario, expected, status
    ("TC-LP-001", "Login", "Login valid (standard_user)", "Masuk ke halaman Products", "PASS"),
    ("TC-LP-004", "Login", "Login locked_out_user", "Tampil pesan 'locked out'", "PASS"),
    ("TC-LP-003", "Login", "Logout via burger menu", "Kembali ke halaman login", "PASS"),
    ("TC-MP-006", "Home", "Sorting Price (low to high)", "Produk terurut harga menaik", "PASS"),
    ("TC-CP-001", "Cart", "Add item ke keranjang", "Item masuk, badge update", "PASS"),
    ("TC-CKO-002", "Checkout", "Isi info checkout valid", "Lanjut ke overview", "PASS"),
    ("TC-CKO-013", "Checkout", "Pesan sukses order", "'Thank you for your order!'", "PASS"),
    ("TC-CP-008", "Cart", "Tambah item lebih dari 1  (Bug #2)", "Qty item bisa bertambah", "FAIL"),
    ("TC-CKO-008", "Checkout", "Input alamat pengiriman  (Bug #9)", "Ada field alamat pengiriman", "FAIL"),
    ("TC-CKO-009", "Checkout", "Pilih metode pembayaran  (Bug #10)", "Ada opsi metode bayar", "FAIL"),
    ("TC-LP-026", "Login", "Account lockout gagal berulang  (Bug #12)", "Akun terkunci sementara", "FAIL"),
    ("TC-MP-013", "Home", "Menu 'All Items'  (Bug #5)", "Navigasi ke katalog", "FAIL"),
    ("TC-F-007", "Footer", "Link ToS & Privacy Policy  (Bug #6)", "Link bisa diklik & berfungsi", "FAIL"),
]
rows = [["TC ID", "Modul", "Skenario", "Expected Result", "Status"]]
rows += [list(k) for k in KEY]
table(s, Inches(0.4), Inches(1.38), Inches(12.55), [0.13, 0.1, 0.34, 0.31, 0.12], rows,
      fontsize=10, row_h=Inches(0.375))
text(s, Inches(0.4), Inches(6.7), Inches(12.4), Inches(0.35),
     [[("7 PASS (alur inti) + 6 FAIL (bug Priority-High & keamanan) — dipilih dari 99 case.", 11, True, NAVY)]])
footer(s, 6)

# ============================================================ 7 EXECUTION RESULT
s = slide(); header(s, "Test Execution Result", "Metrik & status per modul")
metrics = [(str(TOTAL), "Total Case", NAVY), (str(PASS), "PASS", PASSG),
           (str(FAIL), "FAIL", RED), ("82.8%", "Pass Rate", AMBER), ("17", "Bug", RED)]
mw = Inches(2.35); gap = Inches(0.15); x = Inches(0.55)
for val, lab, c in metrics:
    rect(s, x, Inches(1.4), mw, Inches(1.3), c)
    text(s, x, Inches(1.5), mw, Inches(0.85), [[(val, 38, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, Inches(2.28), mw, Inches(0.4), [[(lab, 12, True, WHITE)]], align=PP_ALIGN.CENTER)
    x += mw + gap
# horizontal pass/fail bars per module
text(s, Inches(0.55), Inches(3.0), Inches(6), Inches(0.4), [[("PASS / FAIL per Modul", 14, True, NAVY)]])
maxw = Inches(8.5); base_x = Inches(2.4)
for i, (nm, tot, p, f, ids) in enumerate(MODULES):
    y = Inches(3.5) + i * Inches(0.55)
    text(s, Inches(0.55), y, Inches(1.8), Inches(0.4), [[(nm, 11.5, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    full = int(maxw * (tot / TOTAL) * 1.6)
    pw = int(full * (p / tot))
    fw = full - pw
    rect(s, base_x, y + Inches(0.05), max(pw, 1), Inches(0.32), PASSG)
    if fw > 0:
        rect(s, base_x + pw, y + Inches(0.05), fw, Inches(0.32), RED)
    text(s, base_x + full + Inches(0.1), y, Inches(2.5), Inches(0.4),
         [[(f"{p} pass / {f} fail", 10.5, False, GREY)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)

# ============================================================ 7 BUG OVERVIEW
s = slide(); header(s, "Bug Report — Ringkasan", "17 bug • SauceDemo (Bug Report).xlsx")
# metric cards
bm = [("17", "Total Bug", NAVY), ("17", "Opened", AMBER), ("4", "Severity Medium", AMBER), ("3", "Priority High", RED)]
x = Inches(0.55); mw = Inches(2.95)
for val, lab, c in bm:
    rect(s, x, Inches(1.4), mw, Inches(1.25), c)
    text(s, x, Inches(1.5), mw, Inches(0.85), [[(val, 36, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, Inches(2.28), mw, Inches(0.4), [[(lab, 11.5, True, WHITE)]], align=PP_ALIGN.CENTER)
    x += mw + Inches(0.2)
# severity + priority distribution tables
text(s, Inches(0.55), Inches(3.0), Inches(6), Inches(0.4), [[("Distribusi Severity", 14, True, NAVY)]])
sev = [["Severity", "Jumlah"], ["High", "0"], ["Medium", "4"], ["Low", "13"], ["Critical", "0"]]
table(s, Inches(0.55), Inches(3.45), Inches(5.9), [0.6, 0.4], sev, fontsize=12, row_h=Inches(0.4))
text(s, Inches(6.95), Inches(3.0), Inches(6), Inches(0.4), [[("Distribusi Priority", 14, True, NAVY)]])
pri = [["Priority", "Jumlah"], ["High", "3"], ["Medium", "10"], ["Low", "4"]]
table(s, Inches(6.95), Inches(3.45), Inches(5.9), [0.6, 0.4], pri, fontsize=12, row_h=Inches(0.4))
text(s, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.7),
     [[("Catatan: ", 12, True, NAVY),
       ("Tidak ada bug Critical/High-severity. Semua 17 bug masih ber-status Opened. "
        "3 bug Priority-High menyangkut alur transaksi inti (lihat slide berikut).", 12, False, DARK)]],
     line_spacing=1.15)
footer(s, 8)

# ============================================================ 8 BUG LIST
s = slide(); header(s, "Daftar Bug (17)", "Bug ID • Severity • Priority • Judul")
bugs = [
    ("1", "Low", "Low", "No message when cart is empty"),
    ("2", "Low", "High", "User cannot add an item more than 1"),
    ("3", "Low", "Low", "Product name format wrong (T-Shirt Red)"),
    ("4", "Low", "Low", "Product description format wrong (Backpack)"),
    ("5", "Medium", "Medium", "'All Items' menu label does nothing"),
    ("6", "Medium", "Medium", "ToS & Privacy Policy links not clickable"),
    ("7", "Low", "Low", "Checkout accepts invalid inputs"),
    ("8", "Low", "Medium", "Shipment serial same for multiple orders"),
    ("9", "Medium", "High", "Delivery address field missing in checkout"),
    ("10", "Medium", "High", "Payment method selection missing"),
    ("11", "Low", "Medium", "Checkout allowed with empty cart"),
    ("12-18", "Low", "Medium", "Account not locked after failed logins (6 roles)"),
]
rows = [["ID", "Severity", "Priority", "Judul Bug"]]
rows += [[a, b, c, d] for (a, b, c, d) in bugs]
table(s, Inches(0.5), Inches(1.4), Inches(12.4), [0.1, 0.16, 0.16, 0.58], rows,
      fontsize=10.5, row_h=Inches(0.41))
footer(s, 9)

# ============================================================ 9 CRITICAL BUGS
def bug_card(s, x, y, w, h, bid, title, sev, pri, exp, act):
    rect(s, x, y, w, h, LIGHT)
    rect(s, x, y, w, Inches(0.7), NAVY)
    text(s, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.6),
         [[("BUG #" + bid, 12, True, GREEN)], [(title, 12.5, True, WHITE)]], space_after=1)
    rect(s, x + Inches(0.15), y + Inches(0.82), Inches(1.85), Inches(0.3), AMBER if sev == "Medium" else GREY)
    text(s, x + Inches(0.15), y + Inches(0.82), Inches(1.85), Inches(0.3),
         [[("Severity: " + sev, 10, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x + Inches(2.1), y + Inches(0.82), Inches(1.6), Inches(0.3), RED)
    text(s, x + Inches(2.1), y + Inches(0.82), Inches(1.6), Inches(0.3),
         [[("Priority: " + pri, 10, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.18), y + Inches(1.25), w - Inches(0.36), h - Inches(1.3),
         [[("Expected: ", 10.5, True, PASSG), (exp, 10.5, False, DARK)],
          [("Actual: ", 10.5, True, RED), (act, 10.5, False, DARK)]],
         line_spacing=1.1, space_after=6)

s = slide(); header(s, "Bug Priority-High (3)", "Menyangkut alur transaksi inti")
bug_card(s, Inches(0.4), Inches(1.35), Inches(4.1), Inches(5.3),
    "2", "Tidak bisa menambah item > 1", "Low", "High",
    "Kuantitas item dapat ditambah lebih dari 1.",
    "User tidak dapat menambah lebih dari 1 tiap item.")
bug_card(s, Inches(4.65), Inches(1.35), Inches(4.1), Inches(5.3),
    "9", "Field alamat pengiriman hilang", "Medium", "High",
    "Sistem mengizinkan input & review alamat pengiriman.",
    "Tidak ada input field untuk delivery address.")
bug_card(s, Inches(8.9), Inches(1.35), Inches(4.0), Inches(5.3),
    "10", "Pemilihan metode bayar hilang", "Medium", "High",
    "Sistem mengizinkan pemilihan metode pembayaran.",
    "Tidak ada opsi pemilihan payment method.")
footer(s, 10)

# ============================================================ 10 ANALYSIS
s = slide(); header(s, "Analisis Kualitas Sistem", "Defect kritis • dampak • risiko • retest")
blocks = [
    ("Defect Paling Kritis", "Bug #9 (alamat pengiriman) & #10 (metode bayar) — Priority High. "
        "Untuk produk e-commerce nyata, ketiadaan keduanya membuat transaksi tidak lengkap."),
    ("Dampak ke Pengguna / Bisnis", "Order tak bisa diselesaikan dengan benar (tanpa bayar/alamat). "
        "Bug #2 (qty >1) menurunkan nilai pesanan. Account-lockout absen = risiko keamanan brute-force."),
    ("Risiko jika Tetap Dirilis", "Kehilangan penjualan & data pengiriman, kelemahan keamanan login, "
        "menurunnya kepercayaan. Mayoritas severity Low, namun celah fungsional inti tetap ada."),
    ("Test Case untuk Retesting", "Semua 17 case FAIL setelah perbaikan: prioritas Checkout "
        "(TC-CKO-004/007/008/009/015) & Login lockout (TC-LP-026…031)."),
]
y = Inches(1.4)
for t, d in blocks:
    rect(s, Inches(0.55), y, Inches(12.25), Inches(1.18), LIGHT)
    rect(s, Inches(0.55), y, Inches(0.12), Inches(1.18), GREEN)
    text(s, Inches(0.85), y + Inches(0.1), Inches(11.8), Inches(1.0),
         [[(t, 14, True, NAVY)], [(d, 12, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=1.1)
    y += Inches(1.32)
footer(s, 11)

# ============================================================ 11 RELEASE
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.6), SW, Inches(1.5), AMBER)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.6),
     [[("RELEASE RECOMMENDATION", 18, True, GREEN)]])
text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.1), [[("Keputusan Rilis", 40, True, WHITE)]])
text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.2),
     [[("CONDITIONAL GO", 46, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.9), Inches(4.3), Inches(11.6), Inches(2.6),
     [[("Pass rate 82.8% & tanpa bug Critical — layak rilis SETELAH syarat berikut:", 15, True, WHITE)],
      [("1.  Perbaiki 3 bug Priority-High: #2 (qty), #9 (alamat), #10 (pembayaran).", 14, False, WHITE)],
      [("2.  Perbaiki celah keamanan account-lockout (Bug #12-18).", 14, False, WHITE)],
      [("3.  Bug severity Low lainnya boleh fast-follow patch.", 14, False, WHITE)],
      [("4.  RETEST seluruh 17 case FAIL sebelum go-live.", 14, False, WHITE)]],
     line_spacing=1.25, space_after=6)
footer(s, 12)

# ============================================================ 12 CLOSE
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.35), SW, Pt(4), GREEN)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0), [[("Terima Kasih", 48, True, WHITE)]])
text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6),
     [[("Studi Kasus Pengujian Perangkat Lunak — SauceDemo", 16, False, GREEN)]])
text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.6),
     [[("99 Test Case  •  82 PASS / 17 FAIL  •  17 Bug  •  Keputusan: Conditional GO", 13, False, RGBColor(0xB8,0xC4,0xD6))]])
text(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.5),
     [[("Tim: Zaki & Galih", 12, False, RGBColor(0xB8,0xC4,0xD6))]])

out = "/home/glianalabs/Documents/GitHub/software-testing/SauceDemo_Test_Report.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
